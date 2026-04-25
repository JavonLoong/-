from pptx import Presentation
from pptx.util import Inches, Pt, Emu
import os

prs = Presentation(r'd:\虚拟C盘\腾讯\智能电子产品创新实践-阶段项目汇报PPT-腾讯混元3D (1).pptx')

print(f"Slides count: {len(prs.slides)}")
print(f"Slide width: {prs.slide_width} ({Emu(prs.slide_width).inches:.2f} inches)")
print(f"Slide height: {prs.slide_height} ({Emu(prs.slide_height).inches:.2f} inches)")

# Check slide layouts
for i, layout in enumerate(prs.slide_master.slide_layouts):
    print(f"Layout {i}: {layout.name}")

print("\n" + "="*80)

for i, slide in enumerate(prs.slides):
    print(f"\n{'='*40} Slide {i+1} {'='*40}")
    layout_name = slide.slide_layout.name if slide.slide_layout else "Unknown"
    print(f"Layout: {layout_name}")
    
    for shape in slide.shapes:
        print(f"\n  Shape: {shape.name}")
        print(f"    Type: {shape.shape_type}")
        print(f"    Position: left={shape.left}, top={shape.top}, width={shape.width}, height={shape.height}")
        
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    print(f"    Text: {text[:300]}")
        
        if hasattr(shape, 'image'):
            try:
                img = shape.image
                print(f"    Image: content_type={img.content_type}, size={len(img.blob)} bytes")
                # Save image
                ext = img.content_type.split('/')[-1]
                if ext == 'jpeg':
                    ext = 'jpg'
                img_path = f'd:\\虚拟C盘\\腾讯\\extracted_slide{i+1}_{shape.name}.{ext}'
                with open(img_path, 'wb') as f:
                    f.write(img.blob)
                print(f"    Saved to: {img_path}")
            except Exception as e:
                print(f"    Image error: {e}")
