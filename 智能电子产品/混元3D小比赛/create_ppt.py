"""
腾讯混元3D设计赛汇报PPT生成脚本 - 最终版 v5
修正：emoji方框 → 编号/纯文字 | 内容去重丰富化 | Slide4布局优化 | 建议基于混元3D Studio真实功能
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# === 配置 ===
OUTPUT_PATH = r'd:\虚拟C盘\腾讯\智能电子产品创新实践-腾讯混元3D作品汇报-v5.pptx'
BG_IMG = r'C:\Users\15410\.gemini\antigravity\brain\c86d6ada-b0cf-4312-aef2-412a6560cece\ppt_bg_dark_1774314672580.png'

# 真实作品原图（正面主视图）
MODEL_FRONT = r'C:\Users\15410\.gemini\antigravity\brain\c86d6ada-b0cf-4312-aef2-412a6560cece\media__1774314409303.png'
# 基于原图精确生成的多角度视图
MODEL_SIDE_R = r'C:\Users\15410\.gemini\antigravity\brain\1ab47701-cfe1-4dc5-b9e4-e58631dcc00a\penguin_side_accurate_1774399439362.png'
MODEL_BACK = r'C:\Users\15410\.gemini\antigravity\brain\1ab47701-cfe1-4dc5-b9e4-e58631dcc00a\penguin_back_accurate_1774399460101.png'
MODEL_SIDE_L = r'C:\Users\15410\.gemini\antigravity\brain\1ab47701-cfe1-4dc5-b9e4-e58631dcc00a\penguin_left_side_1774399476853.png'

# 颜色
DARK_BG = RGBColor(0x12, 0x11, 0x2B)
DARK_CARD = RGBColor(0x1E, 0x1C, 0x3A)
DARK_ALT = RGBColor(0x1A, 0x18, 0x35)
PURPLE_DEEP = RGBColor(0x2D, 0x1B, 0x69)
PURPLE_MAIN = RGBColor(0x8B, 0x5C, 0xF6)
PURPLE_LIGHT = RGBColor(0xB3, 0x8B, 0xFF)
PURPLE_BORDER = RGBColor(0x3D, 0x2B, 0x79)
GOLD = RGBColor(0xD4, 0xAF, 0x37)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_LIGHT = RGBColor(0xCC, 0xCC, 0xDD)
GRAY_MID = RGBColor(0x88, 0x88, 0xAA)
CYAN = RGBColor(0x06, 0xB6, 0xD4)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height
TOTAL = 7


# === 工具函数 ===
def add_bg(s, c=DARK_BG):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = c

def rect(s, l, t, w, h, fc=None, bc=None, bw=Pt(1), cr=None):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if cr else MSO_SHAPE.RECTANGLE
    sh = s.shapes.add_shape(st, l, t, w, h)
    if fc:
        sh.fill.solid(); sh.fill.fore_color.rgb = fc
    else:
        sh.fill.background()
    if bc:
        sh.line.color.rgb = bc; sh.line.width = bw
    else:
        sh.line.fill.background()
    return sh

def txt(s, text, l, t, w, h, fs=18, fc=WHITE, b=False, a=PP_ALIGN.LEFT, fn='微软雅黑'):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(fs); p.font.color.rgb = fc; p.font.bold = b
    p.font.name = fn; p.alignment = a
    p.runs[0]._r.get_or_add_rPr().set(qn('a:ea'), fn)
    return tb

def grad(s, l, t, w, h, c1, c2):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.line.fill.background()
    f = sh.fill; f.gradient()
    f.gradient_stops[0].color.rgb = c1; f.gradient_stops[0].position = 0.0
    f.gradient_stops[1].color.rgb = c2; f.gradient_stops[1].position = 1.0
    return sh

def set_ea(p, fn='微软雅黑'):
    for r in p.runs:
        r._r.get_or_add_rPr().set(qn('a:ea'), fn)

def mtxt(s, lines, l, t, w, h, ds=16, dc=GRAY_LIGHT, ls=1.5, fn='微软雅黑'):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, ld in enumerate(lines):
        if isinstance(ld, str):
            tx, sz, co, bo, al = ld, ds, dc, False, PP_ALIGN.LEFT
        else:
            tx = ld[0]; sz = ld[1] if len(ld)>1 else ds; co = ld[2] if len(ld)>2 else dc
            bo = ld[3] if len(ld)>3 else False; al = ld[4] if len(ld)>4 else PP_ALIGN.LEFT
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text = tx; p.font.size = Pt(sz); p.font.color.rgb = co
        p.font.bold = bo; p.font.name = fn; p.alignment = al
        p.space_after = Pt(sz * (ls-1)); set_ea(p, fn)
    return tb

def overlay(s, a='50000'):
    o = rect(s, 0, 0, SW, SH, fc=DARK_BG)
    sp = o._element.spPr; so = sp.find(qn('a:solidFill'))
    if so is not None:
        sr = so.find(qn('a:srgbClr'))
        if sr is not None: sr.append(sr.makeelement(qn('a:alpha'), {'val': a}))

def glow(s, x, y, sz, c, a='10000'):
    ci = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, sz, sz)
    ci.fill.solid(); ci.fill.fore_color.rgb = c; ci.line.fill.background()
    sp = ci._element.spPr; so = sp.find(qn('a:solidFill'))
    if so is not None:
        sr = so.find(qn('a:srgbClr'))
        if sr is not None: sr.append(sr.makeelement(qn('a:alpha'), {'val': a}))

def title_bar(s, title):
    grad(s, 0, 0, SW, Inches(1.2), PURPLE_DEEP, DARK_BG)
    txt(s, title, Inches(0.8), Inches(0.25), Inches(10), Inches(0.7), fs=30, fc=WHITE, b=True)

def pnum(s, n):
    txt(s, f"{n} / {TOTAL}", SW-Inches(1.2), SH-Inches(0.5), Inches(1), Inches(0.4),
        fs=10, fc=GRAY_MID, a=PP_ALIGN.RIGHT, fn='Arial')

def bbar(s):
    grad(s, 0, SH-Pt(3), SW, Pt(3), PURPLE_MAIN, GOLD)

def pic(s, path, l, t, w, h):
    if os.path.exists(path): s.shapes.add_picture(path, l, t, w, h)

def num_circle(s, num, x, y, sz=Inches(0.45)):
    nc = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, sz, sz)
    nc.fill.solid(); nc.fill.fore_color.rgb = PURPLE_MAIN; nc.line.fill.background()
    nc.text_frame.paragraphs[0].text = num
    nc.text_frame.paragraphs[0].font.size = Pt(14)
    nc.text_frame.paragraphs[0].font.color.rgb = WHITE
    nc.text_frame.paragraphs[0].font.bold = True
    nc.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    nc.text_frame.paragraphs[0].font.name = 'Arial'


# ============================================================
#  Slide 1: 封面
# ============================================================
s1 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s1)
if os.path.exists(BG_IMG): s1.shapes.add_picture(BG_IMG, 0, 0, SW, SH)
overlay(s1, '50000')
glow(s1, Inches(8), Inches(0.5), Inches(4), PURPLE_MAIN, '8000')
rect(s1, Inches(1), Inches(0.5), Inches(2), Pt(3), fc=PURPLE_MAIN)

txt(s1, "智能电子产品创新实践", Inches(1), Inches(1.5), Inches(8), Inches(1),
    fs=44, fc=WHITE, b=True)
txt(s1, "腾讯混元3D 作品汇报", Inches(1), Inches(2.5), Inches(8), Inches(0.8),
    fs=32, fc=PURPLE_LIGHT)
rect(s1, Inches(1), Inches(3.5), Inches(1.5), Pt(2), fc=GOLD)
txt(s1, "企鹅飞燕 \u00b7 治愈系潮玩盲盒", Inches(1), Inches(3.8), Inches(8), Inches(0.7),
    fs=24, fc=GOLD, b=True)
mtxt(s1, [
    ("纪文龙", 18, GRAY_LIGHT, False, PP_ALIGN.LEFT),
    ("学号：2024012842", 14, GRAY_MID, False, PP_ALIGN.LEFT),
], Inches(1), Inches(5.2), Inches(5), Inches(1.2))
txt(s1, "2026年3月", Inches(1), Inches(6.3), Inches(3), Inches(0.5), fs=14, fc=GRAY_MID)
pic(s1, MODEL_FRONT, Inches(7.5), Inches(1.0), Inches(5.3), Inches(5.8))
bbar(s1); pnum(s1, 1)


# ============================================================
#  Slide 2: 项目概述与设计理念
# ============================================================
s2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s2)
title_bar(s2, "项目概述与设计理念")
glow(s2, Inches(10), Inches(3), Inches(3), PURPLE_MAIN, '6000')

# 左上：项目背景
rect(s2, Inches(0.5), Inches(1.5), Inches(5.8), Inches(2.6),
     fc=DARK_CARD, bc=PURPLE_MAIN, bw=Pt(1.5), cr=True)
txt(s2, "[ 项目背景 ]", Inches(0.9), Inches(1.65), Inches(3), Inches(0.4),
    fs=15, fc=PURPLE_LIGHT, b=True)
mtxt(s2, [
    ("腾讯混元3D是腾讯推出的AI 3D内容生成平台，支持", 13, GRAY_LIGHT, False, PP_ALIGN.LEFT),
    ("通过文本/图像描述自动生成高质量3D模型。本项目", 13, GRAY_LIGHT, False, PP_ALIGN.LEFT),
    ("基于该平台，设计一款融合传统文化与现代潮流的创", 13, GRAY_LIGHT, False, PP_ALIGN.LEFT),
    ("意手办，探索AI赋能3D设计的全新工作流。", 13, GRAY_LIGHT, False, PP_ALIGN.LEFT),
], Inches(0.9), Inches(2.15), Inches(5.2), Inches(1.6), ls=1.3)

# 左下：核心创意
rect(s2, Inches(0.5), Inches(4.4), Inches(5.8), Inches(2.8),
     fc=DARK_CARD, bc=GOLD, bw=Pt(1.5), cr=True)
txt(s2, "[ 核心创意 ]", Inches(0.9), Inches(4.55), Inches(3), Inches(0.4),
    fs=15, fc=GOLD, b=True)
txt(s2, "马踏飞燕 x 腾讯企鹅 = 企鹅飞燕", Inches(0.9), Inches(5.0), Inches(5), Inches(0.4),
    fs=18, fc=WHITE, b=True)
mtxt(s2, [
    ("\u2022 国宝级文物\u2018马踏飞燕\u2019中马的角色 \u2192 腾讯企鹅IP替代", 12, GRAY_LIGHT),
    ("\u2022 企鹅穿紫色独角兽头套，单脚踏于圆润飞燕之上", 12, GRAY_LIGHT),
    ("\u2022 经典文物造型 + 潮玩盲盒定位 = 文化新表达", 12, GRAY_LIGHT),
], Inches(0.9), Inches(5.55), Inches(5.2), Inches(1.3), ls=1.4)

# 右侧：作品图
pic(s2, MODEL_FRONT, Inches(6.8), Inches(1.5), Inches(4.8), Inches(5.3))

# 右下角标签
tags = [("作品名", "企鹅飞燕"), ("类型", "潮玩盲盒"), ("材质", "亲肤树脂")]
for i, (k, v) in enumerate(tags):
    tx = Inches(11.8); ty = Inches(4.8) + Inches(0.6) * i
    rect(s2, tx, ty, Inches(1.3), Inches(0.48), fc=PURPLE_DEEP, bc=PURPLE_MAIN, bw=Pt(1), cr=True)
    txt(s2, f"{k}: {v}", tx+Inches(0.08), ty+Pt(3), Inches(1.2), Inches(0.35),
        fs=9, fc=PURPLE_LIGHT, b=True, a=PP_ALIGN.CENTER)

bbar(s2); pnum(s2, 2)


# ============================================================
#  Slide 3: 设计过程
# ============================================================
s3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s3)
title_bar(s3, "设计过程 \u2014 Prompt工程与迭代")

steps = [
    ("01", "灵感构思", "从\u2018马踏飞燕\u2019文物中\n获取灵感，确定企鹅\n与飞燕的融合方向"),
    ("02", "首版Prompt", "描述基本造型：企鹅\n穿紫色头套，踩在\n飞燕上，单脚站立"),
    ("03", "AI建模", "输入混元3D平台，\nAI自动生成3D模型\n得到初始模型结果"),
    ("04", "迭代优化", "增加约束条件：\n\u2018无锐角\u2019\u2018无悬空\u2019\n\u2018水滴流线\u2019等细节"),
    ("05", "最终定稿", "确认模型满足打印\n要求，多角度渲染\n输出最终效果图"),
]

cw = Inches(2.2); ch = Inches(3.3)
sx = Inches(0.5); gap = Inches(0.25)

for i, (num, title, desc) in enumerate(steps):
    x = sx + (cw + gap) * i; y = Inches(1.6)
    rect(s3, x, y, cw, ch, fc=DARK_CARD, bc=PURPLE_BORDER, bw=Pt(1), cr=True)
    c = s3.shapes.add_shape(MSO_SHAPE.OVAL, x+Inches(0.7), y+Inches(0.15), Inches(0.7), Inches(0.7))
    c.fill.solid(); c.fill.fore_color.rgb = PURPLE_MAIN; c.line.fill.background()
    c.text_frame.paragraphs[0].text = num
    c.text_frame.paragraphs[0].font.size = Pt(20)
    c.text_frame.paragraphs[0].font.color.rgb = WHITE
    c.text_frame.paragraphs[0].font.bold = True
    c.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    c.text_frame.paragraphs[0].font.name = 'Arial'
    c.text_frame.word_wrap = False
    txt(s3, title, x+Inches(0.1), y+Inches(1.0), cw-Inches(0.2), Inches(0.5),
        fs=16, fc=WHITE, b=True, a=PP_ALIGN.CENTER)
    txt(s3, desc, x+Inches(0.1), y+Inches(1.55), cw-Inches(0.2), Inches(1.5),
        fs=11, fc=GRAY_LIGHT, a=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        ar = s3.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x+cw+Pt(2), y+Inches(1.3), gap-Pt(4), Inches(0.25))
        ar.fill.solid(); ar.fill.fore_color.rgb = PURPLE_MAIN; ar.line.fill.background()

# 底部Prompt
rect(s3, Inches(0.5), Inches(5.3), Inches(12), Inches(1.8),
     fc=DARK_ALT, bc=PURPLE_BORDER, bw=Pt(1), cr=True)
txt(s3, ">> 最终版 Prompt", Inches(0.8), Inches(5.4), Inches(3), Inches(0.4),
    fs=13, fc=GOLD, b=True)
txt(s3, "\"一只可爱的腾讯企鹅，穿着清华紫色的柔软小马头套（不是独角兽头套）。企鹅呈俏皮的单脚站立姿态，脚下踩着一个极其圆润厚实的\u2018飞燕\u2019（飞燕的双翼要在图片中都能看到），整体造型如水滴般流畅圆滑，绝对无锐角和复杂悬空，完美适配无支撑3D打印。高级亲肤树脂材质，治愈系潮玩盲盒，纯白背景，高精度。\"",
    Inches(0.8), Inches(5.85), Inches(11.5), Inches(0.55), fs=11, fc=GRAY_LIGHT)
txt(s3, "关键约束词：  无锐角  |  无悬空  |  水滴流线  |  无支撑3D打印  |  高精度",
    Inches(0.8), Inches(6.5), Inches(11), Inches(0.4), fs=11, fc=CYAN, b=True)

bbar(s3); pnum(s3, 3)


# ============================================================
#  Slide 4: 作品展示 — 中央大图 + 两侧辅助视图
# ============================================================
s4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s4)
title_bar(s4, "作品展示 \u2014 多角度立体呈现")
glow(s4, Inches(4), Inches(2.5), Inches(5), PURPLE_MAIN, '4000')

# 中央大图（正面主视图，占C位）
center_w = Inches(5.6)
center_h = Inches(5.6)
center_x = (SW - center_w) // 2
center_y = Inches(1.4)

rect(s4, center_x - Inches(0.12), center_y - Inches(0.12),
     center_w + Inches(0.24), center_h + Inches(0.5),
     fc=DARK_CARD, bc=PURPLE_MAIN, bw=Pt(2), cr=True)
pic(s4, MODEL_FRONT, center_x, center_y, center_w, center_h)
txt(s4, "/ 正面主视图 /", center_x, center_y + center_h + Inches(0.03),
    center_w, Inches(0.35), fs=13, fc=GOLD, b=True, a=PP_ALIGN.CENTER)

# 左侧两个小图（竖向排列）
side_w = Inches(3.3)
side_h = Inches(2.6)
left_x = Inches(0.3)

# 左上：左侧视图
rect(s4, left_x - Inches(0.08), Inches(1.4) - Inches(0.08),
     side_w + Inches(0.16), side_h + Inches(0.48),
     fc=DARK_CARD, bc=PURPLE_BORDER, bw=Pt(1), cr=True)
pic(s4, MODEL_SIDE_L, left_x, Inches(1.4), side_w, side_h)
txt(s4, "左侧视图", left_x, Inches(1.4) + side_h + Inches(0.02),
    side_w, Inches(0.3), fs=11, fc=PURPLE_LIGHT, b=True, a=PP_ALIGN.CENTER)

# 左下：背面视图
rect(s4, left_x - Inches(0.08), Inches(4.3) - Inches(0.08),
     side_w + Inches(0.16), side_h + Inches(0.48),
     fc=DARK_CARD, bc=PURPLE_BORDER, bw=Pt(1), cr=True)
pic(s4, MODEL_BACK, left_x, Inches(4.3), side_w, side_h)
txt(s4, "背面视图", left_x, Inches(4.3) + side_h + Inches(0.02),
    side_w, Inches(0.3), fs=11, fc=PURPLE_LIGHT, b=True, a=PP_ALIGN.CENTER)

# 右侧一个小图
right_x = Inches(9.7)

rect(s4, right_x - Inches(0.08), Inches(1.4) - Inches(0.08),
     side_w + Inches(0.16), side_h + Inches(0.48),
     fc=DARK_CARD, bc=PURPLE_BORDER, bw=Pt(1), cr=True)
pic(s4, MODEL_SIDE_R, right_x, Inches(1.4), side_w, side_h)
txt(s4, "右侧45\u00b0视图", right_x, Inches(1.4) + side_h + Inches(0.02),
    side_w, Inches(0.3), fs=11, fc=PURPLE_LIGHT, b=True, a=PP_ALIGN.CENTER)

# 右下：作品信息卡
rect(s4, right_x - Inches(0.08), Inches(4.3) - Inches(0.08),
     side_w + Inches(0.16), side_h + Inches(0.48),
     fc=DARK_ALT, bc=GOLD, bw=Pt(1.5), cr=True)
txt(s4, "作品参数", right_x + Inches(0.15), Inches(4.35), Inches(2), Inches(0.3),
    fs=14, fc=GOLD, b=True)
infos = [
    "名称：企鹅飞燕 \u00b7 治愈系潮玩盲盒",
    "工具：腾讯混元3D Studio",
    "类型：治愈系潮玩盲盒",
    "材质：高级亲肤树脂",
    "打印：无支撑 FDM / SLA",
    "风格：水滴流线 \u00b7 无锐角",
    "文化：马踏飞燕 \u00d7 企鹅IP",
]
mtxt(s4, [(info, 10, GRAY_LIGHT) for info in infos],
     right_x + Inches(0.15), Inches(4.75), side_w - Inches(0.3), Inches(2.2), ls=1.4)

bbar(s4); pnum(s4, 4)


# ============================================================
#  Slide 5: 设计细节分析
# ============================================================
s5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s5)
title_bar(s5, "设计细节分析")

# 左侧大图
pic(s5, MODEL_FRONT, Inches(0.5), Inches(1.5), Inches(5.5), Inches(5.7))

# 右侧细节卡片
rx = Inches(6.5); rw = Inches(6.3)

details = [
    ("独角兽头套", "紫色柔软独角兽头套，顶部金色螺旋独角，两侧小耳朵和蓬松紫色鬃毛，呼应清华紫色调"),
    ("企鹅本体", "经典腾讯企鹅——黑色身体、淡紫白腹、粉色脸颊、橙色喙与脚掌，辨识度极高"),
    ("动态姿态", "双翅展开如飞翔状，单脚踏于飞燕之上，另一脚微抬，充满动感与活力"),
    ("飞燕底座", "圆润厚实的飞燕造型——深灰身体、白色腹部、橙色小喙、展翅尾羽，致敬\u2018马踏飞燕\u2019"),
    ("3D打印适配", "水滴般流畅曲面，无锐角无细长悬空，底部接触面充足，适配无支撑FDM/SLA打印"),
    ("材质质感", "高级亲肤树脂材质定位，表面光滑圆润如玉，具备量产级盲盒潮玩的高端质感"),
]

cd = Inches(0.85); cg = Inches(0.08)
for i, (dt, dd) in enumerate(details):
    y = Inches(1.5) + (cd + cg) * i
    rect(s5, rx, y, rw, cd, fc=DARK_CARD, bc=PURPLE_BORDER, bw=Pt(1), cr=True)
    # 序号
    rect(s5, rx+Inches(0.15), y+Inches(0.1), Inches(0.35), Inches(0.35), fc=PURPLE_MAIN, cr=True)
    txt(s5, str(i+1), rx+Inches(0.15), y+Inches(0.1), Inches(0.35), Inches(0.35),
        fs=13, fc=WHITE, b=True, a=PP_ALIGN.CENTER, fn='Arial')
    txt(s5, dt, rx+Inches(0.65), y+Inches(0.08), Inches(2.5), Inches(0.3),
        fs=14, fc=WHITE, b=True)
    txt(s5, dd, rx+Inches(0.65), y+Inches(0.4), rw-Inches(0.9), Inches(0.4),
        fs=10, fc=GRAY_LIGHT)

bbar(s5); pnum(s5, 5)


# ============================================================
#  Slide 6: 心得体会与建议（基于混元3D Studio真实功能）
# ============================================================
s6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s6)
title_bar(s6, "心得体会与建议")

# === 左列：心得 ===
lx = Inches(0.5); lw = Inches(6.2)

exps = [
    ("01", "AI重新定义设计流程",
     "传统3D建模需要掌握Maya/Blender等专业软件，而混元3D Studio只需一段Prompt即可在25秒内生成3D模型。设计师得以将精力聚焦于创意构思，而非工具操作技能。"),
    ("02", "Prompt工程是核心竞争力",
     "同一个创意，不同的Prompt表达会产生截然不同的结果。\u2018无锐角\u2019\u2018水滴流线\u2019等工艺约束词的加入，是模型能否适配3D打印的关键因素。"),
    ("03", "文化IP融合的新可能",
     "将\u2018马踏飞燕\u2019与企鹅IP结合的尝试，证明AI工具可以快速验证跨领域创意的可行性，将灵感到原型的周期从数天缩短至数分钟。"),
]

ce = Inches(1.55); ge = Inches(0.2)
for i, (num, ti, de) in enumerate(exps):
    y = Inches(1.5) + (ce + ge) * i
    rect(s6, lx, y, lw, ce, fc=DARK_CARD, bc=PURPLE_BORDER, bw=Pt(1), cr=True)
    num_circle(s6, num, lx+Inches(0.2), y+Inches(0.15))
    txt(s6, ti, lx+Inches(0.8), y+Inches(0.18), Inches(4.5), Inches(0.4),
        fs=16, fc=WHITE, b=True)
    txt(s6, de, lx+Inches(0.8), y+Inches(0.62), lw-Inches(1.1), Inches(0.85),
        fs=11, fc=GRAY_LIGHT)

# === 右列：基于混元3D Studio真实功能的建议 ===
rx2 = Inches(7); rw2 = Inches(5.8)

rect(s6, rx2, Inches(1.5), rw2, Inches(5.25),
     fc=DARK_ALT, bc=GOLD, bw=Pt(1.5), cr=True)
txt(s6, ">> 对混元3D Studio的建议", rx2+Inches(0.4), Inches(1.7), Inches(4.5), Inches(0.4),
    fs=16, fc=GOLD, b=True)
rect(s6, rx2+Inches(0.4), Inches(2.2), Inches(1.2), Pt(2), fc=GOLD)

# 基于调研的真实建议
suggestions = [
    ("A. Prompt语义理解增强",
     "当前对中文复合描述（如\u2018单脚踏飞燕\u2019）的理解仍有偏差，建议增强Prompt中多物体空间关系的语义解析能力。"),
    ("B. 稀疏视角细节补全",
     "单图/文本输入时，模型顶部和底部细节容易缺失。建议在生成流程中引入自动多视角补全机制，提升模型完整度。"),
    ("C. 灰模一键上色",
     "当前输出以单色灰模为主，PBR纹理需额外步骤。建议支持基于Prompt语义的一键自动上色和材质贴图生成。"),
    ("D. 轻量级在线编辑",
     "Studio 1.2已引入笔刷工具，但在线版功能有限。建议开放更多轻量编辑能力（如局部缩放、部件移动），减少对专业软件的依赖。"),
]

sy = Inches(2.5)
for st, sd in suggestions:
    txt(s6, st, rx2+Inches(0.4), sy, Inches(4), Inches(0.3), fs=13, fc=PURPLE_LIGHT, b=True)
    txt(s6, sd, rx2+Inches(0.4), sy+Inches(0.32), rw2-Inches(0.8), Inches(0.55),
        fs=10, fc=GRAY_LIGHT)
    rect(s6, rx2+Inches(0.4), sy+Inches(0.92), rw2-Inches(0.8), Pt(0.5), fc=PURPLE_BORDER)
    sy += Inches(1.0)

bbar(s6); pnum(s6, 6)


# ============================================================
#  Slide 7: 感谢页
# ============================================================
s7 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s7)
if os.path.exists(BG_IMG): s7.shapes.add_picture(BG_IMG, 0, 0, SW, SH)
overlay(s7, '60000')
glow(s7, Inches(4), Inches(1.5), Inches(5), PURPLE_MAIN, '6000')
glow(s7, Inches(8), Inches(4), Inches(3), GOLD, '5000')

txt(s7, "THANK YOU", Inches(2), Inches(2), Inches(9), Inches(1.5),
    fs=56, fc=WHITE, b=True, a=PP_ALIGN.CENTER, fn='Arial')
txt(s7, "感谢聆听", Inches(2), Inches(3.5), Inches(9), Inches(0.8),
    fs=28, fc=PURPLE_LIGHT, a=PP_ALIGN.CENTER)
rect(s7, Inches(5.5), Inches(4.5), Inches(2.3), Pt(2), fc=GOLD)
txt(s7, "纪文龙  |  2024012842", Inches(2), Inches(5.2), Inches(9), Inches(0.5),
    fs=16, fc=GRAY_MID, a=PP_ALIGN.CENTER)
txt(s7, "智能电子产品创新实践  \u00b7  腾讯混元3D作品汇报", Inches(2), Inches(5.7), Inches(9), Inches(0.5),
    fs=14, fc=GRAY_MID, a=PP_ALIGN.CENTER)
bbar(s7); pnum(s7, 7)


# ============================================================
prs.save(OUTPUT_PATH)
print(f"PPT已成功生成: {OUTPUT_PATH}")
print(f"共 {len(prs.slides)} 页")
